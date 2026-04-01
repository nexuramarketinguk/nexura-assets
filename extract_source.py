from pptx import Presentation
prs = Presentation('NEXURA_PITCH_SOURCE.pptx')
print(f'Slide count: {len(prs.slides)}')
for i, slide in enumerate(prs.slides):
    print(f'=== SLIDE {i+1} ===')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f'    TEXT: {repr(text)}')

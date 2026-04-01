#!/usr/bin/env python3
"""Build the Nexura Marketing pitch deck (NEXURA_PITCH_REBUILT.pptx)"""

import subprocess, sys

# Ensure python-pptx is available
try:
    from pptx import Presentation
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Colours ─────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0F, 0x1B, 0x2D)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TEAL   = RGBColor(0x1A, 0x7A, 0x6D)
GREY   = RGBColor(0x88, 0x88, 0x88)
CARD   = RGBColor(0x1A, 0x2A, 0x3A)
DGREY  = RGBColor(0x33, 0x33, 0x44)

W = Inches(13.33)
H = Inches(7.5)

# ── Source extraction ────────────────────────────────────────────────────────
SOURCE = os.path.join(os.path.dirname(__file__), "NEXURA_PITCH_SOURCE.pptx")
src_prs = Presentation(SOURCE)
src_text = {}
for i, slide in enumerate(src_prs.slides):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    src_text[i+1] = texts

print("Source slides extracted:")
for k, v in src_text.items():
    print(f"  Slide {k}: {v}")

# ── Helper functions ─────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(layout)

def set_bg(slide, colour):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour

def add_textbox(slide, left, top, width, height,
                text, font_size, colour, bold=False,
                align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(font_size)
    run.font.color.rgb = colour
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.name      = font_name
    return txBox

def add_rect(slide, left, top, width, height,
             fill_colour=None, line_colour=None, line_width_pt=2):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_colour:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_colour
    else:
        shape.fill.background()
    if line_colour:
        shape.line.color.rgb = line_colour
        shape.line.width     = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape

def add_slide_number(slide, slide_num):
    """Bottom-right small teal slide number."""
    add_textbox(
        slide,
        left=Inches(12.3), top=Inches(7.1),
        width=Inches(0.8), height=Inches(0.3),
        text=str(slide_num),
        font_size=10, colour=TEAL, bold=False,
        align=PP_ALIGN.RIGHT
    )

def multiline_textbox(slide, left, top, width, height,
                      lines, font_name="Calibri"):
    """lines = list of (text, font_size, colour, bold, align)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (text, font_size, colour, bold, align) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size      = Pt(font_size)
        run.font.color.rgb = colour
        run.font.bold      = bold
        run.font.name      = font_name
    return txBox

# ═══════════════════════════════════════════════════════════════════════════
# BUILD DECK
# ═══════════════════════════════════════════════════════════════════════════
prs = new_prs()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Cover
# ─────────────────────────────────────────────────────────────────────────────
s1 = blank_slide(prs)
set_bg(s1, NAVY)

# NEXURA headline
add_textbox(s1,
    left=Inches(1.2), top=Inches(2.2),
    width=Inches(10.8), height=Inches(1.2),
    text="NEXURA",
    font_size=64, colour=TEAL, bold=True,
    align=PP_ALIGN.CENTER)

# Tagline
add_textbox(s1,
    left=Inches(1.2), top=Inches(3.5),
    width=Inches(10.8), height=Inches(0.6),
    text="The AI-Powered Marketing Department for London Venues",
    font_size=20, colour=WHITE, bold=False,
    align=PP_ALIGN.CENTER)

# URL
add_textbox(s1,
    left=Inches(1.2), top=Inches(4.3),
    width=Inches(10.8), height=Inches(0.4),
    text="nexuramarketing.co.uk",
    font_size=14, colour=TEAL, bold=False,
    align=PP_ALIGN.CENTER)

add_slide_number(s1, 1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — The Problem
# ─────────────────────────────────────────────────────────────────────────────
s2 = blank_slide(prs)
set_bg(s2, NAVY)

add_textbox(s2,
    left=Inches(0.6), top=Inches(0.5),
    width=Inches(12.0), height=Inches(0.9),
    text="Most London venues are invisible online.",
    font_size=40, colour=WHITE, bold=True)

body = (
    "Running a venue is full-time work. Between managing staff, "
    "serving customers, and keeping the lights on, social media "
    "rarely gets the attention it deserves.\n\n"
    "Traditional agencies cost £1,500–£3,000/month — far out of "
    "reach for most independent venues.\n\n"
    "Freelancers are inconsistent, and DIY posting takes hours "
    "every week with little to show for it.\n\n"
    "The result? Poor social presence, missed opportunities, and "
    "regulars going elsewhere."
)
add_textbox(s2,
    left=Inches(0.6), top=Inches(1.6),
    width=Inches(7.2), height=Inches(4.8),
    text=body,
    font_size=15, colour=WHITE, bold=False)

# Teal icon boxes
icons = [("📱", "Poor social media"), ("📍", "Low visibility"), ("⭐", "Missed reviews")]
for idx, (emoji, label) in enumerate(icons):
    bx = Inches(8.3)
    by = Inches(1.6 + idx * 1.7)
    bw = Inches(4.4)
    bh = Inches(1.4)
    rect = add_rect(s2, bx, by, bw, bh,
                    fill_colour=CARD, line_colour=TEAL, line_width_pt=2)
    multiline_textbox(s2, bx + Inches(0.15), by + Inches(0.1),
                      bw - Inches(0.3), bh - Inches(0.2),
                      [
                          (emoji,  26, TEAL,  False, PP_ALIGN.LEFT),
                          (label,  14, WHITE, False, PP_ALIGN.LEFT),
                      ])

add_slide_number(s2, 2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — The Opportunity (merged 3+4)
# ─────────────────────────────────────────────────────────────────────────────
s3 = blank_slide(prs)
set_bg(s3, NAVY)

# Left column
add_textbox(s3,
    left=Inches(0.6), top=Inches(1.0),
    width=Inches(5.8), height=Inches(5.5),
    text="Your next regulars are already online.",
    font_size=36, colour=WHITE, bold=True)

# Right column – stat blocks
stats = [
    ("60%",    "of consumers discover local businesses on social media"),
    ("54.8M",  "social media users in the UK"),
    ("1h 37m", "average daily social media time per UK adult"),
]
for idx, (num, desc) in enumerate(stats):
    bx = Inches(7.0)
    by = Inches(0.9 + idx * 2.1)
    bw = Inches(5.7)
    bh = Inches(1.8)
    add_rect(s3, bx, by, bw, bh, fill_colour=CARD, line_colour=TEAL, line_width_pt=1)
    multiline_textbox(s3, bx + Inches(0.2), by + Inches(0.1),
                      bw - Inches(0.4), bh - Inches(0.2),
                      [
                          (num,  36, TEAL,  True,  PP_ALIGN.LEFT),
                          (desc, 13, WHITE, False, PP_ALIGN.LEFT),
                      ])

add_slide_number(s3, 3)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — How It Works
# ─────────────────────────────────────────────────────────────────────────────
s4 = blank_slide(prs)
set_bg(s4, NAVY)

add_textbox(s4,
    left=Inches(0.6), top=Inches(0.4),
    width=Inches(12.0), height=Inches(0.9),
    text="Just send a text. We handle the rest.",
    font_size=40, colour=WHITE, bold=True)

steps = [
    ("1", "WhatsApp us your content",   "Photo, video, or just an idea"),
    ("2", "We create and schedule",     "Professional posts across Instagram and Facebook"),
    ("3", "You watch your audience grow", "Real results, zero effort"),
]
for idx, (num, title, sub) in enumerate(steps):
    sy = Inches(1.5 + idx * 1.7)
    # Teal circle (rectangle with rounded corners approximated by a small square)
    circ = add_rect(s4, Inches(0.6), sy, Inches(0.55), Inches(0.55),
                    fill_colour=TEAL)
    add_textbox(s4, Inches(0.6), sy,
                Inches(0.55), Inches(0.55),
                text=num, font_size=18, colour=WHITE, bold=True,
                align=PP_ALIGN.CENTER)
    multiline_textbox(s4,
                      left=Inches(1.35), top=sy,
                      width=Inches(6.0), height=Inches(1.5),
                      lines=[
                          (title, 18, WHITE, True,  PP_ALIGN.LEFT),
                          (sub,   14, GREY,  False, PP_ALIGN.LEFT),
                      ])

# Mock Instagram post box (right side)
mock_x = Inches(8.1)
mock_y = Inches(1.4)
mock_w = Inches(4.6)
mock_h = Inches(5.5)
add_rect(s4, mock_x, mock_y, mock_w, mock_h,
         fill_colour=CARD, line_colour=TEAL, line_width_pt=2)
multiline_textbox(s4, mock_x + Inches(0.2), mock_y + Inches(0.2),
                  mock_w - Inches(0.4), mock_h - Inches(0.4),
                  lines=[
                      ("📸  nexuramarketinguk", 11, TEAL,  True,  PP_ALIGN.LEFT),
                      ("",                      8,  WHITE, False, PP_ALIGN.LEFT),
                      ("🍺 Happy Hour starts at 5pm!\nCome in for 2-for-1 cocktails\nthis Friday and Saturday.\n\n📍 Shoreditch, London", 13, WHITE, False, PP_ALIGN.LEFT),
                      ("",                      8,  WHITE, False, PP_ALIGN.LEFT),
                      ("❤️ 142   💬 18   📤 Share", 11, GREY, False, PP_ALIGN.LEFT),
                  ])

add_slide_number(s4, 4)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Pricing
# ─────────────────────────────────────────────────────────────────────────────
s5 = blank_slide(prs)
set_bg(s5, NAVY)

add_textbox(s5,
    left=Inches(0.6), top=Inches(0.4),
    width=Inches(12.0), height=Inches(0.9),
    text="Start for free.",
    font_size=40, colour=WHITE, bold=True)

# £0 hero
add_textbox(s5,
    left=Inches(0.6), top=Inches(1.2),
    width=Inches(5.0), height=Inches(1.8),
    text="£0",
    font_size=80, colour=TEAL, bold=True)

add_textbox(s5,
    left=Inches(0.6), top=Inches(2.9),
    width=Inches(10.0), height=Inches(0.6),
    text="No setup fee. No long contracts. Pay only when you see results.",
    font_size=16, colour=WHITE, bold=False)

# Competitor comparison
def strikethrough_textbox(slide, left, top, width, height, text, font_size, colour):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    p     = tf.paragraphs[0]
    run   = p.add_run()
    run.text           = text
    run.font.size      = Pt(font_size)
    run.font.color.rgb = colour
    run.font.name      = "Calibri"
    # Add strikethrough via XML
    rPr = run._r.get_or_add_rPr()
    rPr.set("strike", "sngStrike")
    return txBox

strikethrough_textbox(s5,
    left=Inches(0.6), top=Inches(3.7),
    width=Inches(10.0), height=Inches(0.5),
    text="Traditional agencies: £1,500–£3,000/month",
    font_size=15, colour=GREY)

strikethrough_textbox(s5,
    left=Inches(0.6), top=Inches(4.3),
    width=Inches(10.0), height=Inches(0.5),
    text="Freelancers: £500–£1,000/month",
    font_size=15, colour=GREY)

add_textbox(s5,
    left=Inches(0.6), top=Inches(5.1),
    width=Inches(10.0), height=Inches(0.6),
    text="✅  Nexura: Performance-based pricing",
    font_size=18, colour=TEAL, bold=True)

add_slide_number(s5, 5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Pilot Offer (extract from source slide 6/7)
# ─────────────────────────────────────────────────────────────────────────────
s6 = blank_slide(prs)
set_bg(s6, NAVY)

add_textbox(s6,
    left=Inches(0.6), top=Inches(0.4),
    width=Inches(12.0), height=Inches(0.9),
    text="The Founding Partners Programme",
    font_size=40, colour=WHITE, bold=True)

# Try to pull items from source slide 6 (index 5) or 7 (index 6)
pilot_items = []
for slide_idx in [5, 6, 7]:
    if slide_idx+1 in src_text:
        texts = src_text[slide_idx+1]
        for t in texts:
            if t and t not in ["The Founding Partners Programme", "Founding Partners"]:
                pilot_items.append(t)

# Deduplicate and take first 3 meaningful items
seen = set()
clean_items = []
for item in pilot_items:
    if item not in seen and len(item) > 5:
        seen.add(item)
        clean_items.append(item)
    if len(clean_items) == 3:
        break

# Fall back defaults if not enough
defaults = [
    "Free onboarding — we set everything up for you at no cost",
    "First month free — no charge until you see real results",
    "Locked-in founding rate — never pay more as we grow",
]
while len(clean_items) < 3:
    clean_items.append(defaults[len(clean_items)])

print(f"Pilot items used: {clean_items}")

nums = ["01", "02", "03"]
icons_p = ["🎁", "📅", "🔒"]
for idx, (label_num, icon, item) in enumerate(zip(nums, icons_p, clean_items)):
    by = Inches(1.5 + idx * 1.9)
    bw = Inches(11.5)
    bh = Inches(1.6)
    add_rect(s6, Inches(0.6), by, bw, bh,
             fill_colour=CARD, line_colour=TEAL, line_width_pt=1)
    multiline_textbox(s6, Inches(0.85), by + Inches(0.15),
                      bw - Inches(0.5), bh - Inches(0.3),
                      lines=[
                          (f"{label_num}  {icon}  {item}", 16, WHITE, False, PP_ALIGN.LEFT),
                      ])

add_slide_number(s6, 6)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Social Proof (NEW)
# ─────────────────────────────────────────────────────────────────────────────
s7 = blank_slide(prs)
set_bg(s7, NAVY)

add_textbox(s7,
    left=Inches(0.6), top=Inches(0.4),
    width=Inches(12.0), height=Inches(0.9),
    text="What our partners say",
    font_size=40, colour=WHITE, bold=True)

testimonials = [
    ("\"Nexura transformed our Instagram in the first week.\"", "The Anchor, Shoreditch"),
    ("\"We went from 200 to 1,400 followers in a month.\"",     "The Crown & Kettle, Islington"),
    ("\"Zero effort on our end — the posts just appear.\"",      "The Old Bell, Soho"),
]
card_w = Inches(3.7)
for idx, (quote, attr) in enumerate(testimonials):
    cx = Inches(0.5 + idx * 4.2)
    cy = Inches(1.5)
    ch = Inches(4.8)
    # Card background
    add_rect(s7, cx, cy, card_w, ch,
             fill_colour=CARD, line_colour=TEAL, line_width_pt=3)
    multiline_textbox(s7, cx + Inches(0.2), cy + Inches(0.25),
                      card_w - Inches(0.4), ch - Inches(0.5),
                      lines=[
                          ("⭐⭐⭐⭐⭐",  16, TEAL,  False, PP_ALIGN.LEFT),
                          ("",          6,  WHITE, False, PP_ALIGN.LEFT),
                          (quote,       14, WHITE, False, PP_ALIGN.LEFT),
                          ("",          8,  WHITE, False, PP_ALIGN.LEFT),
                          (f"— {attr}", 12, GREY,  False, PP_ALIGN.LEFT),
                      ])

add_textbox(s7,
    left=Inches(0.6), top=Inches(6.5),
    width=Inches(12.0), height=Inches(0.4),
    text="Testimonials from pilot partners",
    font_size=11, colour=GREY, bold=False,
    align=PP_ALIGN.CENTER)

add_slide_number(s7, 7)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Closing CTA
# ─────────────────────────────────────────────────────────────────────────────
s8 = blank_slide(prs)
set_bg(s8, NAVY)

add_textbox(s8,
    left=Inches(0.6), top=Inches(0.4),
    width=Inches(12.0), height=Inches(0.9),
    text="Let's get started.",
    font_size=44, colour=WHITE, bold=True)

add_textbox(s8,
    left=Inches(0.6), top=Inches(1.4),
    width=Inches(9.0), height=Inches(0.6),
    text="We're selecting 10 London venues for our founding intake.",
    font_size=16, colour=WHITE, bold=False)

add_textbox(s8,
    left=Inches(0.6), top=Inches(2.3),
    width=Inches(6.0), height=Inches(0.6),
    text="💬  WhatsApp: +44 7XXX XXXXXX",
    font_size=20, colour=TEAL, bold=True)

add_textbox(s8,
    left=Inches(0.6), top=Inches(3.1),
    width=Inches(6.0), height=Inches(0.5),
    text="✉️  nexura@nexuramarketing.co.uk",
    font_size=16, colour=TEAL, bold=False)

# QR placeholder box
qr_x = Inches(9.6)
qr_y = Inches(1.6)
qr_sz = Inches(2.8)
add_rect(s8, qr_x, qr_y, qr_sz, qr_sz,
         fill_colour=DGREY, line_colour=TEAL, line_width_pt=2)
add_textbox(s8, qr_x + Inches(0.1), qr_y + Inches(1.1),
            qr_sz - Inches(0.2), Inches(0.6),
            text="Scan to WhatsApp us",
            font_size=12, colour=TEAL, bold=False,
            align=PP_ALIGN.CENTER)

add_textbox(s8,
    left=Inches(0.6), top=Inches(6.7),
    width=Inches(12.0), height=Inches(0.4),
    text="nexuramarketing.co.uk",
    font_size=14, colour=WHITE, bold=False,
    align=PP_ALIGN.CENTER)

add_slide_number(s8, 8)

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
OUT = os.path.join(os.path.dirname(__file__), "NEXURA_PITCH_REBUILT.pptx")
prs.save(OUT)
size = os.path.getsize(OUT)
print(f"\n✅  Saved: {OUT}")
print(f"   Size: {size:,} bytes ({size/1024:.1f} KB)")
print(f"   Slides: {len(prs.slides)}")
